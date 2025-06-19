from pathlib import Path
from typing import List, Dict, Any
import kss
from loguru import logger
from config.settings import KOREAN_CHUNK_SIZE, KOREAN_CHUNK_OVERLAP
import re
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request
from googleapiclient.discovery import build
import os.path
import pickle

class DocumentProcessor:
    """문서 처리 클래스"""
    
    def __init__(self):
        self.supported_extensions = {
            '.pdf': self._process_pdf,
            '.docx': self._process_docx,
            '.txt': self._process_txt,
            '.xlsx': self._process_xlsx,
            '.pptx': self._process_pptx,
            '.gsheet': self._process_google_sheet,
            '.gslides': self._process_google_slides
        }
        
        # 구글 API 스코프 설정
        self.SCOPES = [
            'https://www.googleapis.com/auth/spreadsheets.readonly',
            'https://www.googleapis.com/auth/presentations.readonly'
        ]
        
        # 구글 API 클라이언트 초기화
        self._init_google_client()
    
    def _init_google_client(self):
        """구글 API 클라이언트 초기화"""
        creds = None
        token_path = Path('config/token.pickle')
        
        # 토큰 파일이 있으면 로드
        if token_path.exists():
            with open(token_path, 'rb') as token:
                creds = pickle.load(token)
        
        # 유효한 인증 정보가 없으면 새로 인증
        if not creds or not creds.valid:
            if creds and creds.expired and creds.refresh_token:
                creds.refresh(Request())
            else:
                flow = InstalledAppFlow.from_client_secrets_file(
                    'config/credentials.json', self.SCOPES)
                creds = flow.run_local_server(port=0)
            
            # 토큰 저장
            with open(token_path, 'wb') as token:
                pickle.dump(creds, token)
        
        # API 서비스 빌드
        self.sheets_service = build('sheets', 'v4', credentials=creds)
        self.slides_service = build('slides', 'v1', credentials=creds)
    
    def process_document(self, file_path: str) -> List[Dict[str, Any]]:
        """문서를 처리하여 청크로 분할"""
        file_path = Path(file_path)
        if not file_path.exists():
            raise FileNotFoundError(f"파일을 찾을 수 없습니다: {file_path}")
            
        extension = file_path.suffix.lower()
        if extension not in self.supported_extensions:
            raise ValueError(f"지원하지 않는 파일 형식입니다: {extension}")
            
        try:
            processor = self.supported_extensions[extension]
            text = processor(file_path)
            chunks = self._split_text(text)
            return chunks
        except Exception as e:
            logger.error(f"문서 처리 중 오류 발생: {str(e)}")
            raise
    
    def _split_text(self, text: str) -> List[Dict[str, Any]]:
        """텍스트를 청크로 분할"""
        sentences = kss.split_sentences(text)
        chunks = []
        current_chunk = []
        current_length = 0
        
        for sentence in sentences:
            sentence_length = len(sentence)
            if current_length + sentence_length > KOREAN_CHUNK_SIZE:
                if current_chunk:
                    chunks.append({
                        'text': ' '.join(current_chunk),
                        'metadata': {
                            'chunk_size': current_length,
                            'num_sentences': len(current_chunk)
                        }
                    })
                current_chunk = [sentence]
                current_length = sentence_length
            else:
                current_chunk.append(sentence)
                current_length += sentence_length
        
        if current_chunk:
            chunks.append({
                'text': ' '.join(current_chunk),
                'metadata': {
                    'chunk_size': current_length,
                    'num_sentences': len(current_chunk)
                }
            })
        
        return chunks
    
    def _process_pdf(self, file_path: Path) -> str:
        """PDF 파일 처리"""
        import fitz
        doc = fitz.open(file_path)
        text = ""
        for page in doc:
            text += page.get_text()
        return text
    
    def _process_docx(self, file_path: Path) -> str:
        """Word 문서 처리"""
        from docx import Document
        doc = Document(file_path)
        return "\n".join([paragraph.text for paragraph in doc.paragraphs])
    
    def _process_txt(self, file_path: Path) -> str:
        """텍스트 파일 처리"""
        return file_path.read_text(encoding='utf-8')
    
    def _process_xlsx(self, file_path: Path) -> str:
        """Excel 파일 처리"""
        import pandas as pd
        df = pd.read_excel(file_path)
        return df.to_string()
    
    def _process_pptx(self, file_path: Path) -> str:
        """PowerPoint 파일 처리"""
        from pptx import Presentation
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    
    def _process_google_sheet(self, file_path: Path) -> str:
        """구글 시트 처리"""
        try:
            # 파일 ID 추출
            with open(file_path, 'r') as f:
                content = f.read()
                file_id = re.search(r'"url": "https://docs.google.com/spreadsheets/d/([^/]+)"', content)
                if not file_id:
                    raise ValueError("구글 시트 ID를 찾을 수 없습니다.")
                file_id = file_id.group(1)
            
            # 시트 데이터 가져오기
            result = self.sheets_service.spreadsheets().values().get(
                spreadsheetId=file_id,
                range='A1:ZZ'
            ).execute()
            
            # 데이터를 텍스트로 변환
            values = result.get('values', [])
            if not values:
                return ""
            
            text_rows = []
            for row in values:
                text_rows.append('\t'.join(str(cell) for cell in row))
            
            return '\n'.join(text_rows)
            
        except Exception as e:
            logger.error(f"구글 시트 처리 중 오류 발생: {str(e)}")
            raise
    
    def _process_google_slides(self, file_path: Path) -> str:
        """구글 슬라이드 처리"""
        try:
            # 파일 ID 추출
            with open(file_path, 'r') as f:
                content = f.read()
                file_id = re.search(r'"url": "https://docs.google.com/presentation/d/([^/]+)"', content)
                if not file_id:
                    raise ValueError("구글 슬라이드 ID를 찾을 수 없습니다.")
                file_id = file_id.group(1)
            
            # 프레젠테이션 데이터 가져오기
            presentation = self.slides_service.presentations().get(
                presentationId=file_id
            ).execute()
            
            # 슬라이드 텍스트 추출
            slides = presentation.get('slides', [])
            text_slides = []
            
            for slide in slides:
                slide_text = []
                for element in slide.get('pageElements', []):
                    if 'shape' in element:
                        shape = element['shape']
                        if 'text' in shape:
                            text_elements = shape['text'].get('textElements', [])
                            for text_element in text_elements:
                                if 'textRun' in text_element:
                                    slide_text.append(text_element['textRun']['content'])
                
                if slide_text:
                    text_slides.append('\n'.join(slide_text))
            
            return '\n\n'.join(text_slides)
            
        except Exception as e:
            logger.error(f"구글 슬라이드 처리 중 오류 발생: {str(e)}")
            raise 